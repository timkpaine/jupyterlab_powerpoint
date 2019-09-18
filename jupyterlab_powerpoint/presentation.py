import nbformat
from pptx import Presentation as _pptx_presentation


class Part(object):
    def __init__(self, placeholder, value=None):
        self._placeholder = placeholder
        self._name = self._placeholder.name
        self._value = value

    def name(self):
        return self._name

    def placeholder(self):
        return self._placeholder


class Slide(object):
    def __init__(self, slide):
        self._slide = slide
        self._slide_layout = slide.slide_layout
        self._name = slide.name
        self._parts = []
        self._part_types = {ps.name: ps for ps in self._slide.placeholders}
        self._placeholder_types = {ps.name: ps for ps in self._slide_layout.placeholders}

    def parts(self):
        return self._parts

    def part_types(self):
        return self._part_types

    def placeholder_types(self):
        return self._placeholder_types

    def name(self):
        return self._name

    def __len__(self):
        return len(self._parts)

    def layout(self):
        return self._slide_layout

    def new_part(self, type):
        if type not in self.part_types():
            print(list(self._part_types.keys()))
            raise ValueError("Part type not found: {}".format(type))
        return self._part_types[type]

    def slide(self):
        return self._slide_layout


class Presentation(object):
    def __init__(self, master_or_presentation=None):
        self._presentation = _pptx_presentation(master_or_presentation)
        self._slides = []
        self._slide_types = {s.name: s for s in self._presentation.slide_master.slide_layouts}

        for slide in self._presentation.slides:
            self._slides.append(Slide(slide))

    def presentation(self):
        return self._presentation

    def slides(self):
        return self._slides

    def slide_types(self):
        return self._slide_types

    def new_slide(self, type):
        if type not in self.slide_types():
            raise ValueError("Slide type not found: {}".format(type))
        s = Slide(self._presentation.slides.add_slide(self._slide_types[type]))
        self._slides.append(s)
        return s

    def __len__(self):
        return len(self._presentation.slides)

    @staticmethod
    def from_notebook(nb, master):
        '''Convert an nbformat.NotebookNode to a Presentation object'''
        pres = Presentation(master)
        slides = []
        for cell in nb.cells:
            if not cell.metadata.get("jupyterlab_powerpoint"):
                continue
            meta = cell.metadata.get("jupyterlab_powerpoint")

            slide_number = int(meta["slide"])
            type = meta["type"]
            placeholder = meta["placeholder"]

            while slide_number > len(slides):
                slides.append(None)
            if slides[slide_number - 1] is None:
                slides[slide_number - 1] = pres.new_slide(type)
            slide = slides[slide_number - 1]
            part = slide.new_part(placeholder)
            if part.has_text_frame:
                part.text_frame.text = cell['source']

        for slide in slides:
            if slide is None:
                raise Exception("Slide gap")
        return pres

    def to_notebook(self):
        '''Convert a presentation object to a nbformat.NotebookNode object'''
        notebook = nbformat.v4.new_notebook()
        for i, slide in enumerate(self.slides()):
            for j, part in enumerate(slide.part_types().values()):
                cell = nbformat.v4.new_raw_cell()
                if part.has_text_frame:
                    cell.source = part.text or " "
                elif part.has_chart:
                    cell.source = "chart"
                elif part.has_table:
                    cell.source = "table"
                cell.metadata['jupyterlab_powerpoint'] = {}
                cell.metadata['jupyterlab_powerpoint']['slide'] = i
                cell.metadata['jupyterlab_powerpoint']['type'] = slide.name() or slide.layout().name
                cell.metadata['jupyterlab_powerpoint']['placeholder'] = part.name
                notebook.cells.append(cell)
        return notebook

    def to_notebook_master(self):
        '''Convert presentation's slidemaster to a nbformat.NotebookNode object'''
        notebook = nbformat.v4.new_notebook()
        for i, slide in enumerate(self.slide_types().values()):
            for j, placeholder in enumerate(slide.placeholders):
                cell = nbformat.v4.new_raw_cell()
                if placeholder.has_text_frame:
                    cell.source = placeholder.text or " "
                elif placeholder.has_chart:
                    cell.source = "chart"
                elif placeholder.has_table:
                    cell.source = "table"
                cell.metadata['jupyterlab_powerpoint'] = {}
                cell.metadata['jupyterlab_powerpoint']['slide'] = i
                cell.metadata['jupyterlab_powerpoint']['type'] = slide.name
                cell.metadata['jupyterlab_powerpoint']['placeholder'] = placeholder.name
                notebook.cells.append(cell)
        return notebook

    def to_presentation(self, filename):
        self._presentation.save(filename)
